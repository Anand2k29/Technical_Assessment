"""Extract text + rough structure (headings that look like chapters/sections) from
uploaded material. Treat all uploaded content as untrusted: it is parsed into
plain text only and never executed or interpreted as instructions."""
import re
from dataclasses import dataclass, field
from pathlib import Path


@dataclass
class Block:
    text: str
    heading: str | None = None
    chapter: str | None = None
    section: str | None = None
    page: int | None = None


HEADING_RE = re.compile(r"^\s*(chapter|unit|section|part)\s+\w+", re.IGNORECASE)


def _classify_heading(line: str, current_chapter: str | None, current_section: str | None) -> tuple[str | None, str | None]:
    stripped = line.strip()
    if not stripped or len(stripped) > 120:
        return current_chapter, current_section
    if re.match(r"^(chapter|unit)\b", stripped, re.IGNORECASE):
        return stripped, None
    if re.match(r"^(section|part)\b", stripped, re.IGNORECASE) or (stripped.isupper() and len(stripped.split()) <= 8):
        return current_chapter, stripped
    return current_chapter, current_section


def parse_pdf(path: Path) -> list[Block]:
    from pypdf import PdfReader
    reader = PdfReader(str(path))
    blocks: list[Block] = []
    chapter = section = None
    for page_num, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        for line in text.split("\n"):
            chapter, section = _classify_heading(line, chapter, section)
        if text.strip():
            blocks.append(Block(text=text, chapter=chapter, section=section, page=page_num))
    return blocks


def parse_docx(path: Path) -> list[Block]:
    from docx import Document as DocxDocument
    doc = DocxDocument(str(path))
    blocks: list[Block] = []
    chapter = section = None
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = (para.style.name or "").lower() if para.style else ""
        if "heading 1" in style or "title" in style:
            chapter, section = text, None
            continue
        if "heading" in style:
            section = text
            continue
        chapter, section = _classify_heading(text, chapter, section)
        blocks.append(Block(text=text, chapter=chapter, section=section))
    return blocks


def parse_pptx(path: Path) -> list[Block]:
    from pptx import Presentation
    prs = Presentation(str(path))
    blocks: list[Block] = []
    for i, slide in enumerate(prs.slides, start=1):
        title = None
        texts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            t = shape.text_frame.text.strip()
            if not t:
                continue
            if shape == slide.shapes.title:
                title = t
            else:
                texts.append(t)
        body = "\n".join(texts)
        if title or body:
            blocks.append(Block(text=f"{title}\n{body}" if title else body, heading=title, chapter=f"Slide {i}", page=i))
    return blocks


def parse_txt(path: Path) -> list[Block]:
    text = path.read_text(encoding="utf-8", errors="ignore")
    blocks: list[Block] = []
    chapter = section = None
    for para in re.split(r"\n\s*\n", text):
        for line in para.split("\n"):
            chapter, section = _classify_heading(line, chapter, section)
        if para.strip():
            blocks.append(Block(text=para, chapter=chapter, section=section))
    return blocks


PARSERS = {".pdf": parse_pdf, ".docx": parse_docx, ".doc": parse_docx, ".pptx": parse_pptx, ".ppt": parse_pptx, ".txt": parse_txt}


def parse_document(path: Path) -> list[Block]:
    ext = path.suffix.lower()
    if ext not in PARSERS:
        raise ValueError(f"Unsupported file type: {ext}")
    return PARSERS[ext](path)
